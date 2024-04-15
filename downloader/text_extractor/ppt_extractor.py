from pptx import Presentation

def ppt_extractor(file_path):
    text = ""
    prs = Presentation(file_path)
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + " "
    return text